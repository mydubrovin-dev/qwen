#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания PowerPoint презентации "Изучаем Луну" для дошкольников.
Презентация содержит яркие фоны, картинки и наглядные элементы.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_moon_presentation():
    # Создаем презентацию
    prs = Presentation()
    
    # Устанавливаем размер слайда 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Цветовая палитра для детей (яркие, насыщенные цвета)
    COLORS = {
        'dark_blue': RGBColor(25, 25, 112),      # MidnightBlue - космос
        'light_blue': RGBColor(135, 206, 250),   # LightSkyBlue
        'yellow': RGBColor(255, 255, 0),         # Яркий желтый
        'orange': RGBColor(255, 165, 0),         # Оранжевый
        'red': RGBColor(255, 69, 0),             # Красный
        'gray': RGBColor(169, 169, 169),         # Серый - Луна
        'white': RGBColor(255, 255, 255),
        'black': RGBColor(0, 0, 0),
        'green': RGBColor(50, 205, 50),          # LimeGreen
        'purple': RGBColor(148, 0, 211),         # Фиолетовый
    }
    
    def add_background(slide, color1, color2=None):
        """Добавляет градиентный или сплошной фон"""
        if color2 is None:
            # Сплошной цвет
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                0, 0, 
                prs.slide_width, 
                prs.slide_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color1
            shape.line.fill.background()
        else:
            # Градиент через два прямоугольника для простоты
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                0, 0, 
                prs.slide_width, 
                prs.slide_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color1
            shape.line.fill.background()
    
    def add_stars(slide, count=30):
        """Добавляет звезды на фон"""
        import random
        random.seed(42)  # Для воспроизводимости
        for _ in range(count):
            x = random.uniform(0, float(prs.slide_width))
            y = random.uniform(0, float(prs.slide_height))
            size = random.uniform(0.05, 0.15)
            
            star = slide.shapes.add_shape(
                MSO_SHAPE.STAR_5_POINT,
                Inches(x),
                Inches(y),
                Inches(size),
                Inches(size)
            )
            star.fill.solid()
            star.fill.fore_color.rgb = COLORS['yellow']
            star.line.fill.background()
    
    def add_moon_illustration(slide, position='right'):
        """Добавляет иллюстрацию Луны"""
        if position == 'right':
            x = prs.slide_width - Inches(4.5)
            y = Inches(1.5)
        else:
            x = Inches(1)
            y = Inches(1.5)
        
        # Основной круг Луны
        moon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x,
            y,
            Inches(4),
            Inches(4)
        )
        moon.fill.solid()
        moon.fill.fore_color.rgb = COLORS['gray']
        moon.line.color.rgb = COLORS['white']
        moon.line.width = Pt(3)
        
        # Кратеры
        crater_positions = [
            (0.5, 0.8, 0.6),
            (2.5, 1.2, 0.5),
            (1.5, 2.5, 0.7),
            (3.0, 2.8, 0.4),
            (0.8, 3.2, 0.5),
        ]
        
        for cx, cy, size in crater_positions:
            crater = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(cx),
                y + Inches(cy),
                Inches(size),
                Inches(size * 0.8)
            )
            crater.fill.solid()
            crater.fill.fore_color.rgb = RGBColor(139, 139, 139)  # DarkGray
            crater.line.fill.background()
    
    def add_astronaut(slide, position='left'):
        """Добавляет силуэт астронавта"""
        if position == 'left':
            x = Inches(0.5)
        else:
            x = prs.slide_width - Inches(2.5)
        
        y = Inches(3.5)
        
        # Тело
        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            Inches(1.2),
            Inches(2)
        )
        body.fill.solid()
        body.fill.fore_color.rgb = COLORS['white']
        body.line.fill.background()
        
        # Шлем
        helmet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.1),
            y - Inches(0.8),
            Inches(1),
            Inches(1)
        )
        helmet.fill.solid()
        helmet.fill.fore_color.rgb = COLORS['white']
        helmet.line.fill.background()
        
        # Визор шлема
        visor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.2),
            y - Inches(0.6),
            Inches(0.7),
            Inches(0.6)
        )
        visor.fill.solid()
        visor.fill.fore_color.rgb = RGBColor(100, 149, 237)  # CornflowerBlue
        visor.line.fill.background()
    
    def add_title_slide():
        """Создает титульный слайд"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Фон - темный космос
        add_background(slide, COLORS['dark_blue'])
        add_stars(slide, 50)
        
        # Добавляем большую Луну
        moon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            prs.slide_width - Inches(5.5),
            Inches(0.5),
            Inches(5),
            Inches(5)
        )
        moon.fill.solid()
        moon.fill.fore_color.rgb = COLORS['gray']
        moon.line.color.rgb = COLORS['white']
        moon.line.width = Pt(4)
        
        # Заголовок
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(2),
            Inches(8),
            Inches(2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "🌙 Изучаем Луну 🌙"
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['yellow']
        title_para.alignment = PP_ALIGN.LEFT
        
        # Подзаголовок
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(3.8),
            Inches(9),
            Inches(2)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = "На Луне всё не так, как на Земле!"
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = COLORS['white']
        subtitle_para.alignment = PP_ALIGN.LEFT
        
        # Добавляем ракету (простой рисунок)
        rocket_body = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(1),
            Inches(4.5),
            Inches(1.5),
            Inches(2.5)
        )
        rocket_body.fill.solid()
        rocket_body.fill.fore_color.rgb = COLORS['white']
        rocket_body.line.fill.background()
        rocket_body.rotation = 180
        
        # Огонь из ракеты
        flame = slide.shapes.add_shape(
            MSO_SHAPE.TEAR,
            Inches(1.35),
            Inches(7),
            Inches(0.6),
            Inches(0.8)
        )
        flame.fill.solid()
        flame.fill.fore_color.rgb = COLORS['orange']
        flame.line.fill.background()
        flame.rotation = 180
    
    def add_content_slide(title, content_points, icon_type='moon', bg_color=None):
        """Создает слайд с контентом"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Фон
        if bg_color is None:
            bg_color = COLORS['dark_blue']
        add_background(slide, bg_color)
        add_stars(slide, 20)
        
        # Заголовок
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.3),
            Inches(12),
            Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['yellow']
        title_para.alignment = PP_ALIGN.LEFT
        
        # Контент
        content_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.8),
            Inches(7.5),
            Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i > 0:
                para = content_frame.add_paragraph()
            else:
                para = content_frame.paragraphs[0]
            
            para.text = "• " + point
            para.font.size = Pt(28)
            para.font.color.rgb = COLORS['white']
            para.space_after = Pt(20)
        
        # Иконка/иллюстрация справа
        if icon_type == 'moon':
            add_moon_illustration(slide, 'right')
        elif icon_type == 'astronaut':
            add_astronaut(slide, 'right')
        elif icon_type == 'thermometer':
            # Термометр
            thermo = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                prs.slide_width - Inches(3),
                Inches(2),
                Inches(0.8),
                Inches(3)
            )
            thermo.fill.solid()
            thermo.fill.fore_color.rgb = COLORS['white']
            thermo.line.fill.background()
            
            # Красная часть (жарко)
            hot = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                prs.slide_width - Inches(2.9),
                Inches(2.1),
                Inches(0.6),
                Inches(1.4)
            )
            hot.fill.solid()
            hot.fill.fore_color.rgb = COLORS['red']
            hot.line.fill.background()
            
            # Синяя часть (холодно)
            cold = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                prs.slide_width - Inches(2.9),
                Inches(4.2),
                Inches(0.6),
                Inches(0.7)
            )
            cold.fill.solid()
            cold.fill.fore_color.rgb = RGBColor(0, 191, 255)  # DeepSkyBlue
            cold.line.fill.background()
        
        elif icon_type == 'water':
            # Капля воды с крестиком
            drop = slide.shapes.add_shape(
                MSO_SHAPE.TEAR,
                prs.slide_width - Inches(3),
                Inches(2),
                Inches(2),
                Inches(2.5)
            )
            drop.fill.solid()
            drop.fill.fore_color.rgb = RGBColor(0, 191, 255)
            drop.line.fill.background()
            drop.rotation = 180
            
            # Крестик
            cross = slide.shapes.add_shape(
                MSO_SHAPE.CROSS,
                prs.slide_width - Inches(2.5),
                Inches(2.5),
                Inches(1),
                Inches(1)
            )
            cross.fill.solid()
            cross.fill.fore_color.rgb = COLORS['red']
            cross.line.color.rgb = COLORS['white']
            cross.line.width = Pt(3)
        
        elif icon_type == 'gravity':
            # Астронавт прыгает
            add_astronaut(slide, 'right')
            # Стрелка вверх
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.UP_ARROW,
                prs.slide_width - Inches(2),
                Inches(1.5),
                Inches(1),
                Inches(2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['green']
            arrow.line.fill.background()
        
        elif icon_type == 'meteor':
            # Метеориты
            for i in range(3):
                meteor = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    prs.slide_width - Inches(3.5) + Inches(i*0.8),
                    Inches(2 + i*0.5),
                    Inches(0.6),
                    Inches(0.5)
                )
                meteor.fill.solid()
                meteor.fill.fore_color.rgb = RGBColor(105, 105, 105)
                meteor.line.color.rgb = COLORS['orange']
                meteor.line.width = Pt(2)
                
                # Хвост метеора
                tail = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_TRIANGLE,
                    prs.slide_width - Inches(3.7) + Inches(i*0.8),
                    Inches(2.2 + i*0.5),
                    Inches(0.4),
                    Inches(0.3)
                )
                tail.fill.solid()
                tail.fill.fore_color.rgb = COLORS['orange']
                tail.line.fill.background()
                tail.rotation = 45
        
        elif icon_type == 'sun':
            # Солнце с лучами
            sun = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                prs.slide_width - Inches(3.5),
                Inches(1.5),
                Inches(2.5),
                Inches(2.5)
            )
            sun.fill.solid()
            sun.fill.fore_color.rgb = COLORS['yellow']
            sun.line.color.rgb = COLORS['orange']
            sun.line.width = Pt(3)
            
            # Лучи
            for angle in range(0, 360, 45):
                ray = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE,
                    prs.slide_width - Inches(2.25),
                    Inches(2.75),
                    Inches(0.4),
                    Inches(1)
                )
                ray.fill.solid()
                ray.fill.fore_color.rgb = COLORS['yellow']
                ray.line.fill.background()
                ray.rotation = angle
    
    # Создаем слайды
    print("Создание титульного слайда...")
    add_title_slide()
    
    print("Создание слайда 1: Нет воздуха...")
    add_content_slide(
        "😮 Нет воздуха",
        [
            "На Луне нельзя дышать!",
            "Нужен специальный костюм",
            "и баллон с кислородом",
            "как у астронавтов 🚀"
        ],
        icon_type='astronaut',
        bg_color=COLORS['dark_blue']
    )
    
    print("Создание слайда 2: Температура...")
    add_content_slide(
        "🌡️ Очень жарко и холодно",
        [
            "Днём: +127°C 🔥",
            "(очень-очень жарко!)",
            "Ночью: -173°C ❄️",
            "(очень-очень холодно!)"
        ],
        icon_type='thermometer',
        bg_color=RGBColor(70, 130, 180)  # SteelBlue
    )
    
    print("Создание слайда 3: Нет воды...")
    add_content_slide(
        "💧 Нет воды",
        [
            "На Луне нет рек и озёр",
            "Нет дождей и снега",
            "Вода только в виде льда",
            "в некоторых кратерах"
        ],
        icon_type='water',
        bg_color=RGBColor(0, 104, 139)  # DarkCyan
    )
    
    print("Создание слайда 4: Слабая гравитация...")
    add_content_slide(
        "🦘 Слабая гравитация",
        [
            "На Луне ты будешь весить",
            "в 6 раз меньше!",
            "Сможешь прыгать очень высоко",
            "как супергерой! 🌟"
        ],
        icon_type='gravity',
        bg_color=RGBColor(75, 0, 130)  # Indigo
    )
    
    print("Создание слайда 5: Метеориты...")
    add_content_slide(
        "☄️ Много метеоритов",
        [
            "С космоса падают маленькие камни",
            "Они называются метеориты",
            "Но не бойся - они маленькие",
            "и падают редко ⭐"
        ],
        icon_type='meteor',
        bg_color=RGBColor(48, 25, 52)  # Темно-фиолетовый
    )
    
    print("Создание слайда 6: Нет атмосферы...")
    add_content_slide(
        "☀️ Нет атмосферы",
        [
            "Нет защиты от Солнца",
            "Солнечные лучи очень опасны",
            "Нужен защитный костюм",
            "и шлем с фильтром 🕶️"
        ],
        icon_type='sun',
        bg_color=RGBColor(25, 25, 112)  # MidnightBlue
    )
    
    # Сохраняем презентацию
    output_file = "Изучаем_Луну_презентация.pptx"
    prs.save(output_file)
    print(f"\n✅ Презентация успешно создана: {output_file}")
    print(f"Количество слайдов: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    create_moon_presentation()
